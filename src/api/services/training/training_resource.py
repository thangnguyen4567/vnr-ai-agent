import os
from pathlib import Path

import pytesseract
import requests
from langchain_community.document_loaders import (
    Docx2txtLoader,
    PyPDFLoader,
    UnstructuredURLLoader,
)
from langchain_core.documents import Document
from langchain_openai import ChatOpenAI
from pdf2image import convert_from_path
from pptx import Presentation

from src.api.services.training.base_training import Training
from src.config import settings
from src.utils.read_file import generate_random_string


class TrainingResource(Training):
    # Giới hạn số ký tự đọc để tóm tắt. Tài liệu quá lớn chỉ đọc phần đầu
    # trong phạm vi này để tránh vượt context/chi phí LLM. Có thể tinh chỉnh.
    MAX_SUMMARY_CHARS = 15000

    def __init__(self):
        super().__init__()
        self.columns = ["title", "content", "coursemoduleid", "courseid", "source", "coursename"]
        self.model = ChatOpenAI(
            model=settings.LLM_CONFIG["openai"]["model"],
            temperature=settings.LLM_CONFIG["openai"]["temperature"],
            api_key=settings.LLM_CONFIG["openai"]["api_key"],
            base_url=settings.LLM_CONFIG["openai"]["base_url"],
            max_tokens=1500,
        )

    def save_training_data(self, data):

        path = Path(data["source"])
        typefile = path.suffix.lower()
        collection = self.get_collection_name(data)

        try:

            # 1. Đọc toàn bộ text của tài liệu (pdf/docx/pptx/url), có OCR cho PDF ảnh
            full_text = self._extract_text(data["source"], typefile)

            metadata = {}
            for key, value in data.items():
                if key in self.columns:
                    metadata[key] = value

            # 2. Tóm tắt tài liệu bằng AI (làm trước để nếu lỗi thì record cũ vẫn còn)
            summary = self._summarize(full_text, metadata)

            content = "Tài liệu: " + metadata["title"]
            content += "Thuộc lớp học: " + metadata["coursename"]
            content += summary

            document = Document(page_content=content, metadata=metadata)

            # 3. Nếu tài liệu đã tồn tại (cùng coursemoduleid) thì xóa record cũ để
            #    cập nhật lại; chưa có thì tạo mới. Mỗi tài liệu chỉ giữ 1 record.
            for key in self.redis_client.scan_iter(self._doc_key_pattern(collection)):
                coursemoduleid = self.redis_client.hget(key, "coursemoduleid")
                if coursemoduleid is not None and coursemoduleid.decode() == metadata["coursemoduleid"]:
                    self.redis_client.delete(key)

            # 4. Ghi record mới (ghi đè nếu đã tồn tại)
            self.vector_db.add_documents([document], collection)

            self.response["error"] = False
            self.response["message"] = "Training thành công"

            return self.response

        except Exception as e:

            print(str(e))

            self.response["error"] = True
            self.response["message"] = f"Training thất bại: {str(e)}"

            return self.response

    def _extract_text(self, source, typefile):
        """Đọc toàn bộ nội dung text của tài liệu theo từng loại file."""

        if typefile == ".pdf":
            docs = PyPDFLoader(source)
            text = ""
            for doc in docs.load():
                text += doc.page_content

            if text != "":
                return text

            # Xử lý cho case file pdf toàn là hình ảnh -> OCR
            response = requests.get(source)
            name = generate_random_string() + ".pdf"
            with open(name, "wb") as file:
                file.write(response.content)

            images = convert_from_path(name)
            extracted_text = ""
            for image in images:
                extracted_text += pytesseract.image_to_string(image, lang="vie") + "\n"

            os.remove(name)
            return extracted_text

        elif typefile == ".docx":
            response = requests.get(source)
            name = generate_random_string() + ".docx"
            with open(name, "wb") as file:
                file.write(response.content)

            docs = Docx2txtLoader(name).load()
            os.remove(name)
            return "".join(doc.page_content for doc in docs)

        elif typefile == ".pptx":
            response = requests.get(source)
            name = generate_random_string() + ".pptx"
            with open(name, "wb") as file:
                file.write(response.content)

            full_text = ""
            presentation = Presentation(name)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"

            os.remove(name)
            return full_text

        else:
            docs = UnstructuredURLLoader(urls=[source]).load()
            return "".join(doc.page_content for doc in docs)

    def _summarize(self, text, metadata):
        """Dùng AI tóm tắt tài liệu. Chỉ đọc trong phạm vi MAX_SUMMARY_CHARS."""

        text = (text or "").strip()
        if text == "":
            return ""

        # Giới hạn phạm vi đọc để tóm tắt với tài liệu quá lớn
        text = text[: self.MAX_SUMMARY_CHARS]

        system_prompt = (
            "Bạn là trợ lý tóm tắt tài liệu đào tạo. "
            "Hãy tóm tắt nội dung tài liệu bên dưới bằng tiếng Việt, "
            "bao quát đầy đủ các ý chính, khái niệm và từ khóa quan trọng "
            "để phục vụ tìm kiếm ngữ nghĩa. "
            "Chỉ trả về nội dung tóm tắt, không thêm lời dẫn."
        )
        human_prompt = (
            f"Tên tài liệu: {metadata.get('title', '')}\n"
            f"Thuộc lớp học: {metadata.get('coursename', '')}\n\n"
            f"Nội dung tài liệu:\n{text}"
        )

        result = self.model.invoke(
            [
                ("system", system_prompt),
                ("human", human_prompt),
            ]
        )
        return result.content

    def delete_training_data(self, data):

        collection = self.get_collection_name(data)

        try:
            coursemoduleids_set = set(data["coursemoduleids"])
            keys_to_delete = []

            for key in self.redis_client.scan_iter(self._doc_key_pattern(collection)):
                coursemoduleid_redis = self.redis_client.hget(key, "coursemoduleid")
                if coursemoduleid_redis is not None and coursemoduleid_redis.decode() in coursemoduleids_set:
                    keys_to_delete.append(key)

            if keys_to_delete:
                self.redis_client.delete(*keys_to_delete)

            self.response["error"] = False
            self.response["message"] = "Xóa dữ liệu thành công"

            return self.response

        except Exception as e:

            print(str(e))
            self.response["error"] = True
            self.response["message"] = f"Xóa dữ liệu thất bại: {str(e)}"

            return self.response

    def check_training_duplication(self):
        pass

    def get_collection_name(self, data):
        return "resource_" + data["contextdata"]["collection"]
