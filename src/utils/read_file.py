import os
import requests
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, UnstructuredURLLoader, UnstructuredExcelLoader

def generate_random_string(length: int = 8):
    import random, string
    return ''.join(random.choices(string.ascii_lowercase + string.digits, k=length))

def load_file_to_text(filepath: str) -> str:
    """
    Load nhiều loại file (pdf, docx, pptx, url khác) thành string cho LangChain.
    - filepath: đường dẫn hoặc URL
    """
    if filepath is None:
        return "empty"
    
    _, ext = os.path.splitext(filepath)
    typefile = ext.lower()

    text_documents = ""
    try:
        if typefile == '.pdf':
            # Load PDF dạng text
            docs = PyPDFLoader(filepath)
            text = ''.join([doc.page_content.strip() for doc in docs.load()])

            text_documents = text

        elif typefile == '.docx':
            # DOCX
            response = requests.get(filepath, verify=False)
            filename = f"{generate_random_string()}.docx"
            with open(filename, "wb") as f:
                f.write(response.content)

            docs = Docx2txtLoader(filename).load()
            text_documents = ''.join([doc.page_content.strip() for doc in docs])
            os.remove(filename)

        elif typefile == '.pptx':
            # PPTX
            response = requests.get(filepath, verify=False)
            filename = f"{generate_random_string()}.pptx"
            with open(filename, "wb") as f:
                f.write(response.content)

            full_text = ""
            presentation = Presentation(filename)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"

            text_documents = full_text
            os.remove(filename)

        elif typefile == '.xlsx' or typefile == '.xls':
            # XLSX
            response = requests.get(filepath, verify=False)
            filename = f"{generate_random_string()}.xlsx"
            with open(filename, "wb") as f:
                f.write(response.content)
            docs = UnstructuredExcelLoader(filename).load()
            text_documents = ''.join([doc.page_content.strip() for doc in docs])
            os.remove(filename)

        else:
            # URL hoặc định dạng khác
            docs = UnstructuredURLLoader(urls=[filepath]).load()

            for doc in docs:
                text_documents += doc.page_content.strip() + "\n"

        return text_documents

    except Exception as e:
        print(f"❌ Lỗi khi load file {filepath}: {e}")
        return "Lỗi khi load file"
